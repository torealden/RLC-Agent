"""
Rebuild Slide 6 ("The Monday Morning I Finally Fixed") with native PowerPoint
shapes that match the deck's dark theme.  Replaces the PNG screenshot version.

Operates on the v3 file and saves as v4.
"""

import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

SRC = 'conferences/ai_is_not_coming_for_your_job_presentation_v3.pptx'
DST = 'conferences/ai_is_not_coming_for_your_job_presentation_v4.pptx'

# ── Deck palette ───────────────────────────────────────────────────
GOLD       = RGBColor(0xC8, 0x96, 0x3E)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE  = RGBColor(0xEA, 0xE6, 0xE0)
GRAY       = RGBColor(0x6B, 0x65, 0x60)
DARK_BG    = RGBColor(0x1A, 0x19, 0x17)
PANEL_BG   = RGBColor(0x24, 0x22, 0x20)   # slightly lighter than bg for panels
LIGHT      = RGBColor(0xD4, 0xCF, 0xC5)

# Task-block colors (BEFORE)
C_DOWNLOAD = RGBColor(0xC5, 0x5A, 0x11)   # burnt orange
C_ENTRY    = RGBColor(0xB0, 0x30, 0x30)   # dark red
C_FORMAT   = RGBColor(0xBF, 0x8F, 0x00)   # amber
C_READING  = RGBColor(0x7F, 0x60, 0x00)   # dark brown
C_ANALYSIS = RGBColor(0x2E, 0x75, 0xB6)   # blue

# AFTER colors
C_OVERNIGHT = RGBColor(0x2D, 0x55, 0x3D)  # muted forest green
C_AUTO      = RGBColor(0x3D, 0x6B, 0x4F)  # medium green
C_KG        = RGBColor(0x54, 0x82, 0x35)  # bright green

# Summary colors
RED_SUMMARY   = RGBColor(0x8B, 0x3A, 0x3A)
GREEN_SUMMARY = RGBColor(0x3D, 0x6B, 0x4F)


def _in(val):
    """Shorthand for Inches()."""
    return Inches(val)


def add_rounded_rect(slide, left, top, width, height, fill_color,
                     corner_radius=Inches(0.06)):
    """Add a rounded rectangle with solid fill, no border."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # Set corner radius via XML
    sp = shape._element
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    prstGeom = sp.find('.//a:prstGeom', nsmap)
    if prstGeom is not None:
        avLst = prstGeom.find('a:avLst', nsmap)
        if avLst is None:
            avLst = etree.SubElement(prstGeom, '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
        # Clear existing
        for child in list(avLst):
            avLst.remove(child)
        gd = etree.SubElement(avLst, '{http://schemas.openxmlformats.org/drawingml/2006/main}gd')
        gd.set('name', 'adj')
        # Value in 1/50000 of shape dimension; ~8000 gives a nice gentle radius
        gd.set('fmla', f'val {int(corner_radius / min(width, height) * 50000)}')
    return shape


def add_text_box(slide, left, top, width, height, text, size_pt=11,
                 bold=False, italic=False, color=WHITE, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP):
    """Add a text box with formatting."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_multiline_text_box(slide, left, top, width, height, lines, align=PP_ALIGN.CENTER,
                           anchor=MSO_ANCHOR.MIDDLE):
    """Add text box with multiple lines, each a (text, size, bold, italic, color) tuple."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass

    for i, (text, size, bold, italic, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color

    return tb


def add_task_block(slide, left, top, width, block_h, time_h,
                   fill_color, lines, time_text):
    """
    Add a task block: rounded rectangle + centered label text + time underneath.
    lines: list of (text, size_pt, bold, italic, color)
    """
    add_rounded_rect(slide, left, top, width, block_h, fill_color)
    add_multiline_text_box(
        slide, left, top, width, block_h, lines,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
    )
    if time_text:
        add_text_box(
            slide, left, top + block_h + _in(0.02), width, time_h,
            time_text, size_pt=7, italic=True, color=GRAY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP
        )


# ══════════════════════════════════════════════════════════════════
#  Main
# ══════════════════════════════════════════════════════════════════

prs = Presentation(SRC)
slide = prs.slides[5]  # Slide 6 (0-indexed)
slide_w = prs.slide_width   # 9144000
slide_h = prs.slide_height  # 5143500

# ── Clear everything on the slide ──
for shape in list(slide.shapes):
    sp = shape._element
    sp.getparent().remove(sp)

# ── Full-slide dark background ──
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_BG
bg.line.fill.background()

# ── Top accent bar ──
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, _in(0.05))
bar.fill.solid()
bar.fill.fore_color.rgb = GOLD
bar.line.fill.background()

# ── Slide Title ──
add_text_box(
    slide, _in(0.55), _in(0.18), _in(8.5), _in(0.45),
    'The Monday Morning I Finally Fixed',
    size_pt=22, bold=True, color=WHITE, align=PP_ALIGN.LEFT
)

# ════════════════════════════════════════════════════════════════
#  BEFORE section
# ════════════════════════════════════════════════════════════════

before_y = _in(0.72)

# "BEFORE" badge
add_rounded_rect(slide, _in(0.55), before_y, _in(0.85), _in(0.2), RED_SUMMARY)
add_text_box(
    slide, _in(0.55), before_y, _in(0.85), _in(0.2),
    'BEFORE', size_pt=8, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
)

# Subtitle
add_text_box(
    slide, _in(1.55), before_y, _in(4.5), _in(0.2),
    'Manual Data Collection & Entry',
    size_pt=10, color=GRAY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE
)

# ── BEFORE task blocks ──
# Proportional widths based on time spent (total ~205 min mapped to ~9.0" usable)
block_y = _in(1.05)
block_h = _in(0.72)
time_h  = _in(0.18)
gap     = _in(0.06)
x_start = _in(0.35)

# (label_lines, width_inches, color, time_text)
before_tasks = [
    ([('Download', 9, False, False, WHITE),
      ('WASDE PDF', 9, False, False, WHITE)],
     0.92, C_DOWNLOAD, '~25m'),

    ([('Key In', 9, False, False, WHITE),
      ('S&D Numbers', 8, False, False, WHITE)],
     1.05, C_ENTRY, '~30m'),

    ([('Pull CFTC', 8, False, False, WHITE),
      ('from Website', 7.5, False, False, WHITE)],
     0.72, C_DOWNLOAD, '~15m'),

    ([('Update COT', 8.5, False, False, WHITE),
      ('Spreadsheet', 8, False, False, WHITE)],
     0.85, C_ENTRY, '~20m'),

    ([('EIA Ethanol', 8, False, False, WHITE),
      ('Download', 8, False, False, WHITE)],
     0.72, C_DOWNLOAD, '~15m'),

    ([('Format &', 8, False, False, WHITE),
      ('Copy Data', 8, False, False, WHITE)],
     0.72, C_FORMAT, '~15m'),

    ([('Export Sales', 8.5, False, False, WHITE),
      ('Report', 8.5, False, False, WHITE)],
     0.85, C_DOWNLOAD, '~20m'),

    ([('Read Weather', 8, False, False, WHITE),
      ('Emails', 8, False, False, WHITE)],
     0.72, C_READING, '~15m'),

    ([('Crop', 8, False, False, WHITE),
      ('Progress', 7.5, False, False, WHITE)],
     0.55, C_DOWNLOAD, '~10m'),

    ([('Finally:', 9, False, False, WHITE),
      ('Analysis', 9, True, False, WHITE)],
     1.42, C_ANALYSIS, '~40m'),
]

x = x_start
for i, (lines, w_in, color, time_text) in enumerate(before_tasks):
    w = _in(w_in)

    # Dashed separator before Analysis block
    if i == len(before_tasks) - 1:
        sep_x = x - gap // 2
        sep = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, sep_x, block_y - _in(0.04),
            _in(0.02), block_h + _in(0.08)
        )
        sep.fill.solid()
        sep.fill.fore_color.rgb = GOLD
        sep.fill.fore_color.brightness = 0.3
        sep.line.fill.background()

    add_task_block(slide, x, block_y, w, block_h, time_h, color, lines, time_text)
    x += w + gap

# "~20% of morning" annotation pointing to Analysis block
analysis_left = x - _in(1.42) - gap
add_text_box(
    slide, analysis_left - _in(0.1), block_y - _in(0.3), _in(1.6), _in(0.28),
    '~20% of morning\nspent on analysis',
    size_pt=7.5, italic=True, color=GOLD,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM
)

# ── BEFORE summary bar ──
summary_y = block_y + block_h + time_h + _in(0.12)
# Subtle dark panel behind summary
add_rounded_rect(
    slide, _in(0.35), summary_y, _in(9.0), _in(0.26),
    RGBColor(0x30, 0x20, 0x18)
)
add_text_box(
    slide, _in(0.35), summary_y, _in(9.0), _in(0.26),
    '~3 hours of mechanical data work  \u00b7  before any real analysis begins',
    size_pt=10.5, bold=True, color=RED_SUMMARY,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
)

# ════════════════════════════════════════════════════════════════
#  Thin horizontal divider
# ════════════════════════════════════════════════════════════════
div_y = summary_y + _in(0.38)
divider = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, _in(1.5), div_y, _in(6.7), _in(0.015)
)
divider.fill.solid()
divider.fill.fore_color.rgb = RGBColor(0x3A, 0x36, 0x32)
divider.line.fill.background()

# ════════════════════════════════════════════════════════════════
#  AFTER section
# ════════════════════════════════════════════════════════════════

after_label_y = div_y + _in(0.12)

# "AFTER" badge
add_rounded_rect(slide, _in(0.55), after_label_y, _in(0.85), _in(0.2), GREEN_SUMMARY)
add_text_box(
    slide, _in(0.55), after_label_y, _in(0.85), _in(0.2),
    'AFTER', size_pt=8, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
)

# Subtitle
add_text_box(
    slide, _in(1.55), after_label_y, _in(4.5), _in(0.2),
    'AI-Powered Automated Pipeline',
    size_pt=10, color=GRAY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE
)

# "The part that actually matters" annotation
add_text_box(
    slide, _in(5.0), after_label_y, _in(4.5), _in(0.2),
    'The part that actually matters \u2192',
    size_pt=8, bold=True, italic=True, color=GOLD,
    align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE
)

# ── AFTER task blocks ──
after_block_y = after_label_y + _in(0.32)
after_block_h = _in(0.72)

# Block 1: "Ran Overnight" — wide muted green
ran_left = _in(0.35)
ran_w = _in(2.5)
add_rounded_rect(slide, ran_left, after_block_y, ran_w, after_block_h, C_OVERNIGHT)
add_multiline_text_box(
    slide, ran_left, after_block_y, ran_w, after_block_h,
    [('Ran Overnight', 11, True, False, WHITE),
     ('30+ collectors finished', 8.5, False, False, RGBColor(0xB4, 0xD6, 0x9C))],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
)
add_text_box(
    slide, ran_left, after_block_y + after_block_h + _in(0.02), ran_w, time_h,
    'while you slept', size_pt=7, italic=True, color=GRAY,
    align=PP_ALIGN.CENTER
)

# Block 2: "Review Briefing"
rb_left = ran_left + ran_w + gap
rb_w = _in(0.9)
add_task_block(
    slide, rb_left, after_block_y, rb_w, after_block_h, time_h,
    C_AUTO,
    [('Review', 8.5, False, False, WHITE),
     ('Briefing', 8.5, False, False, WHITE)],
    '~5 min'
)

# Block 3: "KG Loaded"
kg_left = rb_left + rb_w + gap
kg_w = _in(0.8)
add_task_block(
    slide, kg_left, after_block_y, kg_w, after_block_h, time_h,
    C_KG,
    [('KG', 9, False, False, WHITE),
     ('Loaded', 8.5, False, False, WHITE)],
    'instant'
)

# Dashed separator before big Analysis block
sep2_x = kg_left + kg_w + gap // 2
sep2 = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, sep2_x, after_block_y - _in(0.04),
    _in(0.02), after_block_h + _in(0.08)
)
sep2.fill.solid()
sep2.fill.fore_color.rgb = GOLD
sep2.fill.fore_color.brightness = 0.3
sep2.line.fill.background()

# Block 4: "Analysis, Insight & Client Communication" — big blue
an_left = kg_left + kg_w + gap
an_w = _in(9.35) - an_left  # Fill to right margin
add_rounded_rect(slide, an_left, after_block_y, an_w, after_block_h, C_ANALYSIS)
add_multiline_text_box(
    slide, an_left, after_block_y, an_w, after_block_h,
    [('Analysis, Insight & Client Communication', 13, True, False, WHITE),
     ('Compare forecasts  \u00b7  Generate reports  \u00b7  Write commentary  \u00b7  Identify opportunities',
      8, False, False, RGBColor(0xCA, 0xDC, 0xFC))],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
)

# ── AFTER summary bar ──
after_summary_y = after_block_y + after_block_h + time_h + _in(0.12)
add_rounded_rect(
    slide, _in(0.35), after_summary_y, _in(9.0), _in(0.26),
    RGBColor(0x18, 0x2C, 0x1E)
)
add_text_box(
    slide, _in(0.35), after_summary_y, _in(9.0), _in(0.26),
    '~80% of morning spent on analysis & insight  \u00b7  data collection happened while you slept',
    size_pt=10.5, bold=True, color=GREEN_SUMMARY,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
)

# ── Footer ──
add_text_box(
    slide, _in(0.55), _in(5.2), _in(1.0), _in(0.22),
    '6', size_pt=9, color=GRAY, align=PP_ALIGN.LEFT
)
add_text_box(
    slide, _in(3.0), _in(5.2), _in(6.0), _in(0.22),
    'Round Lakes Companies  \u00b7  roundlakescommodities.com',
    size_pt=9, color=GRAY, align=PP_ALIGN.RIGHT
)

# ── Save ──
prs.save(DST)
print(f'Saved rebuilt slide 6 to {DST}')
