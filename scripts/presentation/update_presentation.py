"""
Update ai_is_not_coming_for_your_job_presentation_v2.pptx with all improvements.

Changes:
  1. Slide 1:  "Lines of Code Written" → "Lines of Code Written by Hand"
  2. Slide 6:  Replace dense before/after with the generated PNG image
  3. Slide 11: Replace RINs IF→THEN with a grain-accessible causal chain
  4. Slide 12: Insert rule_in_action_timeline.png into the placeholder area
  5. Slide 14: Reframe the 18-month claim to honest aspiration
  6. Slide 22: Soften step 3 of the CTA
  7. Slide 23: Add "Built with Claude Code by Anthropic" credit
  8. NEW slide after 20: "What Surprised Me"
  9. Insert CFTC chart into slide 13 (auto reports demo)
  10. Insert data freshness dashboard into slide 10 (data collection demo)
  11. Fix slide numbers throughout

Saves as _v3.pptx
"""

import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy
from lxml import etree
import os

SRC = 'conferences/ai_is_not_coming_for_your_job_presentation_v2.pptx'
DST = 'conferences/ai_is_not_coming_for_your_job_presentation_v3.pptx'

CHARTS_DIR = 'data/generated_graphics/charts'

# Colors matching the deck
GOLD    = RGBColor(0xC8, 0x96, 0x3E)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0x6B, 0x65, 0x60)
DARK_BG = RGBColor(0x1A, 0x19, 0x17)
LIGHT   = RGBColor(0xD4, 0xCF, 0xC5)
GREEN   = RGBColor(0x3D, 0x6B, 0x4F)
RED     = RGBColor(0x8B, 0x3A, 0x3A)
DARK_TEXT = RGBColor(0x0F, 0x0E, 0x0C)


def find_shape_by_text(slide, search_text):
    """Find first shape containing search_text."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            full = shape.text_frame.text
            if search_text in full:
                return shape
    return None


def find_shapes_by_text(slide, search_text):
    """Find all shapes containing search_text."""
    results = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            if search_text in shape.text_frame.text:
                results.append(shape)
    return results


def set_shape_text(shape, text, size_pt=11, bold=False, italic=False,
                   color=WHITE, alignment=PP_ALIGN.LEFT):
    """Clear shape text and set new text with formatting."""
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_run_to_para(para, text, size_pt=11, bold=False, italic=False, color=WHITE):
    """Add a formatted run to a paragraph."""
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return run


def add_text_box(slide, left, top, width, height, text, size_pt=11,
                 bold=False, italic=False, color=WHITE, alignment=PP_ALIGN.LEFT):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_image_to_slide(slide, image_path, left, top, width=None, height=None):
    """Add an image to a slide."""
    kwargs = {'left': left, 'top': top}
    if width:
        kwargs['width'] = width
    if height:
        kwargs['height'] = height
    return slide.shapes.add_picture(image_path, **kwargs)


def duplicate_slide(prs, slide_index):
    """Duplicate a slide layout and create a new blank slide after slide_index."""
    source = prs.slides[slide_index]
    layout = source.slide_layout

    # Add new slide at end first
    new_slide = prs.slides.add_slide(layout)

    # Clear default placeholders
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    return new_slide


def move_slide(prs, old_index, new_index):
    """Move a slide from old_index to new_index."""
    # Access the sldIdLst from the presentation XML element
    prs_elem = prs.part._element
    nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    sldIdLst = prs_elem.find('.//p:sldIdLst', nsmap)
    slides = list(sldIdLst)
    el = slides[old_index]
    sldIdLst.remove(el)
    target_slides = list(sldIdLst)
    if new_index >= len(target_slides):
        sldIdLst.append(el)
    else:
        ref = target_slides[new_index]
        ref.addprevious(el)


def add_dark_bg_rect(slide, prs_width, prs_height):
    """Add a dark background rectangle covering the full slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs_width, prs_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BG
    shape.line.fill.background()
    # Send to back by moving XML element to front of shape tree
    sp_tree = slide.shapes._spTree
    sp_tree.insert(2, shape._element)  # After the required elements
    return shape


# ══════════════════════════════════════════════════════════════════
#  Main
# ══════════════════════════════════════════════════════════════════

prs = Presentation(SRC)
slides = prs.slides
slide_w = prs.slide_width
slide_h = prs.slide_height

print(f'Loaded {len(slides)} slides, {slide_w}x{slide_h}')

# ── 1. Slide 1: "Lines of Code Written" → "Lines of Code Written by Hand" ──
print('1. Fixing slide 1 — "by Hand"')
s1 = slides[0]
shape = find_shape_by_text(s1, 'Lines of Code Written')
if shape:
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if 'Lines of Code Written' in run.text:
                run.text = run.text.replace('Lines of Code Written',
                                           'Lines of Code Written by Hand')
                print('   Done.')

# ── 2. Slide 6: Replace dense before/after with PNG ──
print('2. Fixing slide 6 — before/after image')
s6 = slides[5]

# Remove all the tiny task block shapes (keep just the title and the top bar)
shapes_to_keep_texts = ['The Monday Morning I Finally Fixed']
shapes_to_remove = []
for shape in s6.shapes:
    keep = False
    # Keep the title
    if shape.has_text_frame:
        txt = shape.text_frame.text.strip()
        if txt in shapes_to_keep_texts:
            keep = True
        # Keep the top accent bar (no text, at top)
    # Keep the first shape (accent bar at y=0)
    if shape.top is not None and shape.top < Emu(100000) and not shape.has_text_frame:
        keep = True
    if not keep:
        shapes_to_remove.append(shape)

# Actually, let's be more surgical - remove everything except title and top bar,
# then add the image
for shape in shapes_to_remove:
    sp = shape._element
    sp.getparent().remove(sp)

# Add the before/after workflow image
ba_path = os.path.join(CHARTS_DIR, 'before_after_workflow.png')
if os.path.exists(ba_path):
    # Position below title, centered
    img_left = Inches(0.3)
    img_top = Inches(1.0)
    img_width = Inches(9.1)
    add_image_to_slide(s6, ba_path, img_left, img_top, width=img_width)
    print('   Added before/after workflow image.')

# ── 3. Slide 11: Replace third IF→THEN rule ──
print('3. Fixing slide 11 — grain-accessible IF→THEN')
s11 = slides[10]
shape = find_shape_by_text(s11, 'RINs stack vs. feedstock')
if shape:
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if 'RINs stack vs. feedstock' in run.text:
                run.text = run.text.replace(
                    'RINs stack vs. feedstock margin divergence \u2192 palm oil export surge warning',
                    'Late safrinha planting in Brazil \u2192 smaller corn crop \u2192 US exports gain share \u2192 bullish Dec corn'
                )
                print('   Done.')

# ── 4. Slide 12: Insert rule_in_action chart ──
print('4. Fixing slide 12 — insert timeline chart')
s12 = slides[11]
# Find and remove the placeholder text about the chart
placeholder = find_shape_by_text(s12, '[ Chart:')
if placeholder:
    sp = placeholder._element
    sp.getparent().remove(sp)

# Also remove the second placeholder line if separate
placeholder2 = find_shape_by_text(s12, 'RINs/LCFS credit-stack anomaly')
if placeholder2 and 'flag vs. Wednesday' in placeholder2.text_frame.text:
    sp = placeholder2._element
    sp.getparent().remove(sp)

ria_path = os.path.join(CHARTS_DIR, 'rule_in_action_timeline.png')
if os.path.exists(ria_path):
    img_left = Inches(0.4)
    img_top = Inches(0.9)
    img_width = Inches(9.0)
    add_image_to_slide(s12, ria_path, img_left, img_top, width=img_width)
    print('   Added rule in action timeline.')

# ── 5. Slide 14: Reframe 18-month claim ──
print('5. Fixing slide 14 — reframe forecast claim')
s14 = slides[13]
shape = find_shape_by_text(s14, '18 months')
if shape:
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if '18 months' in run.text:
                run.text = ('The framework is designed so that over time, '
                           'both the human and the AI improve \u2014 '
                           'symbiotic accuracy tracking that no spreadsheet '
                           'alone can deliver.')
                print('   Done.')

# ── 6. Slide 22: Soften step 3 ──
print('6. Fixing slide 22 — soften CTA step 3')
s22 = slides[21]
shape = find_shape_by_text(s22, 'Encode your first rule')
if shape:
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if 'Encode your first rule' in run.text:
                run.text = ('3.  Write down the one rule you\u2019d tell a new '
                           'hire on day one \u2014 that\u2019s your first '
                           'knowledge graph entry.')
                print('   Done.')

# ── 7. Slide 23: Add Claude Code credit ──
print('7. Fixing slide 23 — add Claude Code credit')
s23 = slides[22]
add_text_box(
    s23,
    left=Inches(1.5), top=Inches(4.5),
    width=Inches(7.0), height=Inches(0.4),
    text='Built with Claude Code by Anthropic',
    size_pt=10, italic=True, color=GRAY,
    alignment=PP_ALIGN.CENTER
)
print('   Done.')

# ── 8. NEW: "What Surprised Me" slide after slide 20 ──
print('8. Creating "What Surprised Me" slide')

# Create new slide using the same layout as slide 20 (Five Lessons)
new_slide = prs.slides.add_slide(slides[19].slide_layout)

# Clear default placeholders
for shape in list(new_slide.shapes):
    sp = shape._element
    sp.getparent().remove(sp)

# Dark background
add_dark_bg_rect(new_slide, slide_w, slide_h)

# Top accent bar
accent_bar = new_slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, 0, slide_w, Emu(54864)
)
accent_bar.fill.solid()
accent_bar.fill.fore_color.rgb = GOLD
accent_bar.line.fill.background()

# Title
add_text_box(
    new_slide,
    left=Inches(0.6), top=Inches(0.35),
    width=Inches(8.5), height=Inches(0.6),
    text='What Surprised Me',
    size_pt=28, bold=True, color=DARK_TEXT,
    alignment=PP_ALIGN.LEFT
)

# Surprise items — three numbered items with explanations
items = [
    ('1', 'The AI remembered my rules better than I did',
     'I\u2019d encoded a crop condition rule six months ago and forgotten about it. '
     'The system flagged it before I even thought to check. Your knowledge compounds '
     'even when you\u2019re not paying attention.'),
    ('2', 'Explaining what I wanted was the hardest part',
     'I\u2019d done these workflows on autopilot for 20 years. Translating intuition '
     'into explicit instructions forced me to think clearly about processes I\u2019d '
     'never articulated \u2014 and made me a better analyst in the process.'),
    ('3', 'The system found signals I\u2019d been missing',
     'Cross-market correlations across 30+ data sources simultaneously \u2014 '
     'no human can hold that many variables in their head. The AI doesn\u2019t get '
     'tired, distracted, or anchored to yesterday\u2019s narrative.'),
]

y_pos = Inches(1.15)
for num, title, desc in items:
    # Number
    add_text_box(
        new_slide,
        left=Inches(0.6), top=y_pos,
        width=Inches(0.6), height=Inches(0.55),
        text=num, size_pt=22, bold=True, color=GOLD,
        alignment=PP_ALIGN.CENTER
    )
    # Title
    add_text_box(
        new_slide,
        left=Inches(1.3), top=y_pos,
        width=Inches(7.8), height=Inches(0.45),
        text=title, size_pt=13, bold=True, color=DARK_TEXT,
        alignment=PP_ALIGN.LEFT
    )
    # Description
    add_text_box(
        new_slide,
        left=Inches(1.3), top=y_pos + Inches(0.42),
        width=Inches(7.8), height=Inches(0.55),
        text=desc, size_pt=11, italic=False, color=GRAY,
        alignment=PP_ALIGN.LEFT
    )
    y_pos += Inches(1.15)

# Footer
add_text_box(
    new_slide,
    left=Inches(0.6), top=Inches(4.65),
    width=Inches(3.0), height=Inches(0.25),
    text='20', size_pt=9, color=GRAY,
    alignment=PP_ALIGN.LEFT
)
add_text_box(
    new_slide,
    left=Inches(2.5), top=Inches(4.65),
    width=Inches(6.5), height=Inches(0.25),
    text='Round Lakes Companies  \u00b7  roundlakescommodities.com',
    size_pt=9, color=GRAY,
    alignment=PP_ALIGN.RIGHT
)

# Move new slide from end (index -1) to after slide 20 (index 20)
total = len(prs.slides)
move_slide(prs, total - 1, 20)
print('   Created and positioned "What Surprised Me" slide.')

# ── 9. Insert CFTC chart into slide 13 (auto reports) ──
print('9. Adding CFTC chart to slide 13')
s13 = slides[12]
cftc_path = os.path.join(CHARTS_DIR, 'cftc_corn_positioning_presentation.png')
if os.path.exists(cftc_path):
    # Add as a smaller inset in the lower portion of the slide
    add_image_to_slide(
        s13, cftc_path,
        left=Inches(4.8), top=Inches(1.2),
        width=Inches(4.8)
    )
    print('   Done.')

# ── 10. Insert data freshness dashboard into slide 10 ──
print('10. Adding data freshness dashboard to slide 10')
s10 = slides[9]
df_path = os.path.join(CHARTS_DIR, 'data_freshness_dashboard_presentation.png')
if os.path.exists(df_path):
    add_image_to_slide(
        s10, df_path,
        left=Inches(4.5), top=Inches(0.8),
        width=Inches(5.0)
    )
    print('   Done.')

# ── 11. Fix slide numbers ──
print('11. Fixing slide numbers throughout')
# After inserting the new slide, total is now 24 slides
# We need to update footer page numbers
# The convention seems to be: some slides have a small text box with the number
for idx, slide in enumerate(prs.slides):
    slide_num = idx + 1
    # Skip title slide (1), section dividers (3, 7, 9, 15, 19)
    # Look for shapes that contain just a small number
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            # Check if this is a page number (1-2 digit number, small font)
            if txt.isdigit() and len(txt) <= 2:
                # Check if it's in the footer area (bottom of slide)
                if shape.top > slide_h * 0.85:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip().isdigit():
                                run.text = str(slide_num)

print('   Done.')

# ── Save ──
prs.save(DST)
print(f'\nSaved updated presentation to {DST}')
print(f'Total slides: {len(prs.slides)}')
