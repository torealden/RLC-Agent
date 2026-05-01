"""
Pass 1: Walk a folder, extract content excerpts, write inventory CSV.

No LLM calls. Just walks the tree, pulls text from PDF/DOCX/XLSX/PPTX/TXT/CSV,
and dumps to CSV with one row per file. Resumable — re-running picks up where
it stopped (files already in the CSV are skipped).

Usage:
    python scripts/dscan_inventory.py "D:/Seagate Recovery/Recoverit 2022-10-13 at 01.36.12"
    python scripts/dscan_inventory.py <folder> --out docs/reference/inv_recovery.csv
    python scripts/dscan_inventory.py <folder> --max-bytes 50_000_000  (skip larger files)

Output CSV columns:
    rel_path, name, ext, size_bytes, mtime, looks_lost,
    extract_status, content_excerpt
"""

from __future__ import annotations

import argparse
import csv
import os
import re
import sys
import traceback
from datetime import datetime
from pathlib import Path
from typing import Optional

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

# --- excerpt config -----------------------------------------------------------
EXCERPT_CHARS = 1500          # how much text to keep per file
MAX_BYTES_DEFAULT = 50_000_000  # skip files larger than 50 MB by default
TEXT_EXTS = {".txt", ".csv", ".log", ".md", ".tsv", ".xml", ".json", ".html", ".htm"}

# --- "looks_lost" heuristic ---------------------------------------------------
# Patterns that suggest the original filename was lost and Recoverit assigned
# a placeholder. Real names like "AYP.accdb" or "Pork.csv" are SHORT but
# meaningful, so length alone isn't enough — we use specific patterns.
LOST_PATTERNS = [
    re.compile(r"^[0-9]{6,}$"),                      # all-numeric (recovery cluster IDs)
    re.compile(r"^~\$"),                              # Office lock/temp files
    re.compile(r"^(document|untitled|new|copy|~|temp|tmp)\d*$", re.I),
    re.compile(r"^(file|recovered|lost|unknown)\d*$", re.I),
    re.compile(r"^[0-9a-f]{16,}$", re.I),             # long hex hashes
]


def looks_lost(stem: str) -> bool:
    return any(p.match(stem) for p in LOST_PATTERNS)


# --- text extractors ----------------------------------------------------------
def excerpt_pdf(path: Path) -> str:
    from pypdf import PdfReader
    r = PdfReader(str(path), strict=False)
    out = []
    for page in r.pages[:3]:                   # only first 3 pages
        try:
            out.append(page.extract_text() or "")
        except Exception:
            continue
        if sum(len(s) for s in out) > EXCERPT_CHARS:
            break
    return "\n".join(out)


def excerpt_docx(path: Path) -> str:
    import docx
    d = docx.Document(str(path))
    out = []
    for p in d.paragraphs:
        out.append(p.text)
        if sum(len(s) for s in out) > EXCERPT_CHARS:
            break
    return "\n".join(out)


def excerpt_xlsx(path: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    bits = [f"sheets={wb.sheetnames}"]
    for sn in wb.sheetnames[:3]:
        ws = wb[sn]
        bits.append(f"\n[{sn}]")
        for r_idx, row in enumerate(ws.iter_rows(values_only=True)):
            cells = [str(c) for c in row[:8] if c is not None][:8]
            if cells:
                bits.append(" | ".join(cells))
            if r_idx > 12:
                break
        if sum(len(b) for b in bits) > EXCERPT_CHARS:
            break
    wb.close()
    return "\n".join(bits)


def excerpt_pptx(path: Path) -> str:
    from pptx import Presentation
    p = Presentation(str(path))
    out = []
    for i, slide in enumerate(p.slides[:8]):
        out.append(f"--- slide {i+1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            out.append(run.text.strip())
        if sum(len(s) for s in out) > EXCERPT_CHARS:
            break
    return "\n".join(out)


def excerpt_text(path: Path) -> str:
    for enc in ("utf-8", "cp1252", "latin-1"):
        try:
            with open(path, "r", encoding=enc, errors="replace") as f:
                return f.read(EXCERPT_CHARS * 2)
        except Exception:
            continue
    return ""


EXTRACTORS = {
    ".pdf":  excerpt_pdf,
    ".docx": excerpt_docx,
    ".xlsx": excerpt_xlsx,
    ".xlsm": excerpt_xlsx,
    ".pptx": excerpt_pptx,
}


def extract_excerpt(path: Path, ext: str) -> tuple[str, str]:
    """Returns (status, excerpt). Status is 'ok', 'skip:reason', or 'err:...'"""
    if ext in EXTRACTORS:
        try:
            text = EXTRACTORS[ext](path)
        except Exception as e:
            return f"err:{type(e).__name__}", ""
    elif ext in TEXT_EXTS:
        try:
            text = excerpt_text(path)
        except Exception as e:
            return f"err:{type(e).__name__}", ""
    else:
        return "skip:not_extractable", ""

    text = (text or "").strip()
    if not text:
        return "ok:empty", ""
    # Collapse internal whitespace, cap length
    text = re.sub(r"\s+", " ", text)[:EXCERPT_CHARS]
    return "ok", text


# --- main walk ---------------------------------------------------------------
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("root", type=Path)
    ap.add_argument("--out", type=Path, default=None,
                    help="Output CSV (default: docs/reference/dscan_<root_slug>.csv)")
    ap.add_argument("--max-bytes", type=int, default=MAX_BYTES_DEFAULT,
                    help="Skip files larger than this (default 50 MB)")
    ap.add_argument("--limit", type=int, default=0, help="Stop after N files (debug)")
    args = ap.parse_args()

    if not args.root.exists():
        sys.exit(f"Path not found: {args.root}")

    if args.out is None:
        slug = re.sub(r"[^A-Za-z0-9_-]+", "_", str(args.root.name)).strip("_")[:60] or "scan"
        args.out = ROOT / "docs" / "reference" / f"dscan_{slug}.csv"
    args.out.parent.mkdir(parents=True, exist_ok=True)

    # Resume support: load already-processed paths
    seen: set[str] = set()
    if args.out.exists():
        with open(args.out, "r", encoding="utf-8", newline="") as f:
            for row in csv.DictReader(f):
                seen.add(row["rel_path"])
        print(f"Resume: {len(seen):,} files already in {args.out.name}")

    fields = ["rel_path", "name", "ext", "size_bytes", "mtime",
              "looks_lost", "extract_status", "content_excerpt"]

    # Open in append mode; write header only if new file
    write_header = not args.out.exists()
    f_out = open(args.out, "a", encoding="utf-8", newline="")
    writer = csv.DictWriter(f_out, fieldnames=fields, quoting=csv.QUOTE_ALL)
    if write_header:
        writer.writeheader()

    n_total = n_done = n_skipped = n_extracted = n_err = 0
    last_progress = datetime.now()

    try:
        for dp, dirs, files in os.walk(args.root):
            for fn in files:
                n_total += 1
                p = Path(dp) / fn
                rel = str(p.relative_to(args.root)).replace("\\", "/")
                if rel in seen:
                    continue

                try:
                    st = p.stat()
                    size = st.st_size
                    mtime = datetime.fromtimestamp(st.st_mtime).isoformat(timespec="seconds")
                except OSError:
                    size = 0
                    mtime = ""

                ext = p.suffix.lower()
                row = {
                    "rel_path": rel,
                    "name": fn,
                    "ext": ext,
                    "size_bytes": size,
                    "mtime": mtime,
                    "looks_lost": looks_lost(p.stem),
                    "extract_status": "",
                    "content_excerpt": "",
                }

                if size > args.max_bytes:
                    row["extract_status"] = "skip:too_large"
                    n_skipped += 1
                else:
                    status, excerpt = extract_excerpt(p, ext)
                    row["extract_status"] = status
                    row["content_excerpt"] = excerpt
                    if status == "ok":
                        n_extracted += 1
                    elif status.startswith("err:"):
                        n_err += 1
                    else:
                        n_skipped += 1

                writer.writerow(row)
                n_done += 1

                # Progress every 5 seconds
                now = datetime.now()
                if (now - last_progress).total_seconds() > 5:
                    print(f"  {n_done:>7,} processed | "
                          f"extracted={n_extracted:,} skipped={n_skipped:,} err={n_err:,}", flush=True)
                    last_progress = now
                    f_out.flush()

                if args.limit and n_done >= args.limit:
                    print(f"Hit --limit {args.limit}, stopping.")
                    return
    except KeyboardInterrupt:
        print("\nInterrupted — partial progress saved. Re-run to resume.")
    finally:
        f_out.close()

    print(f"\nDone. {n_done:,} new files written to {args.out}")
    print(f"  extracted: {n_extracted:,}  skipped: {n_skipped:,}  errors: {n_err:,}")


if __name__ == "__main__":
    main()
